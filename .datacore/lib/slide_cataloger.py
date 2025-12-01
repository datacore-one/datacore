#!/usr/bin/env python3
"""
Slide Cataloger - Extract and index presentations for searchable catalog.

Creates a JSON catalog that enables full-text search while preserving
links to original files for design retrieval.

Usage:
    python slide_cataloger.py <source_dir> <output_dir>

Example:
    python slide_cataloger.py 1-teamspace/4-outbox/archive/presentations 1-teamspace/presentations
"""

import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from collections import Counter
import hashlib

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    import pdfplumber
except ImportError:
    print("Error: pdfplumber not installed. Run: pip install pdfplumber")
    sys.exit(1)


def extract_keywords(text: str, top_n: int = 10) -> list[str]:
    """Extract top keywords from text content."""
    # Common stop words to filter out
    stop_words = {
        'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
        'of', 'with', 'by', 'from', 'as', 'is', 'was', 'are', 'were', 'been',
        'be', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would',
        'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'need',
        'it', 'its', 'this', 'that', 'these', 'those', 'i', 'you', 'we',
        'they', 'he', 'she', 'your', 'our', 'their', 'my', 'his', 'her',
        'what', 'which', 'who', 'whom', 'how', 'when', 'where', 'why',
        'all', 'each', 'every', 'both', 'few', 'more', 'most', 'other',
        'some', 'such', 'no', 'nor', 'not', 'only', 'own', 'same', 'so',
        'than', 'too', 'very', 'just', 'also', 'now', 'here', 'there',
        'then', 'once', 'if', 'because', 'while', 'although', 'though',
        'after', 'before', 'until', 'unless', 'about', 'into', 'through',
        'during', 'above', 'below', 'between', 'under', 'again', 'further',
        'any', 'etc', 'new', 'one', 'two', 'first', 'get', 'use', 'using'
    }

    # Extract words
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())

    # Filter and count
    filtered = [w for w in words if w not in stop_words]
    counts = Counter(filtered)

    # Return top keywords
    return [word for word, _ in counts.most_common(top_n)]


def infer_audience(filename: str, content: str) -> str:
    """Infer target audience from filename and content."""
    filename_lower = filename.lower()
    content_lower = content.lower()

    # Check filename patterns
    if any(x in filename_lower for x in ['investor', 'pitch', 'funding', 'seed', 'series']):
        return 'investor'
    if any(x in filename_lower for x in ['partner', 'collab', 'b2b']):
        return 'partner'
    if any(x in filename_lower for x in ['conf', 'summit', 'meetup', 'talk', 'keynote', 'presentation']):
        return 'conference'
    if any(x in filename_lower for x in ['internal', 'team', 'offsite', 'strategy']):
        return 'internal'

    # Check content patterns
    if any(x in content_lower for x in ['invest', 'valuation', 'cap table', 'runway', 'series']):
        return 'investor'
    if any(x in content_lower for x in ['partnership', 'integrate', 'api access']):
        return 'partner'

    return 'general'


def parse_date_from_filename(filename: str) -> str:
    """Try to extract date from filename."""
    # Try various date patterns
    patterns = [
        r'(\d{4})-(\d{2})-(\d{2})',  # 2018-06-15
        r'(\d{4})(\d{2})(\d{2})',     # 20180615
        r'(\d{2})-(\d{2})-(\d{4})',   # 15-06-2018
        r'(\d{4})-(\d{2})',           # 2018-06
        r'(\d{4})_(\d{2})',           # 2018_06
    ]

    for pattern in patterns:
        match = re.search(pattern, filename)
        if match:
            groups = match.groups()
            if len(groups) == 3:
                if len(groups[0]) == 4:
                    return f"{groups[0]}-{groups[1]}-{groups[2]}"
                else:
                    return f"{groups[2]}-{groups[1]}-{groups[0]}"
            elif len(groups) == 2:
                return f"{groups[0]}-{groups[1]}-01"

    return ""


def extract_pptx(filepath: Path) -> dict | None:
    """Extract content from PPTX file."""
    try:
        prs = Presentation(str(filepath))
    except (PackageNotFoundError, Exception) as e:
        print(f"  Error opening {filepath.name}: {e}")
        return None

    slides = []
    all_text = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        title = ""

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_text.append(text)

                # First text shape is often the title
                if not title and len(text) < 100:
                    title = text

        content = "\n".join(slide_text)
        all_text.append(content)

        if content:  # Only add non-empty slides
            slides.append({
                "number": i,
                "title": title[:100] if title else f"Slide {i}",
                "content": content,
                "keywords": extract_keywords(content, 5)
            })

    full_content = "\n\n".join(all_text)

    return {
        "slide_count": len(prs.slides),
        "slides": slides,
        "full_content": full_content,
        "audience": infer_audience(filepath.name, full_content)
    }


def extract_pdf(filepath: Path) -> dict | None:
    """Extract content from PDF file."""
    try:
        pdf = pdfplumber.open(str(filepath))
    except Exception as e:
        print(f"  Error opening {filepath.name}: {e}")
        return None

    slides = []
    all_text = []

    try:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            all_text.append(text)

            if text.strip():
                lines = text.strip().split('\n')
                title = lines[0][:100] if lines else f"Page {i}"

                slides.append({
                    "number": i,
                    "title": title,
                    "content": text.strip(),
                    "keywords": extract_keywords(text, 5)
                })

        full_content = "\n\n".join(all_text)

        return {
            "slide_count": len(pdf.pages),
            "slides": slides,
            "full_content": full_content,
            "audience": infer_audience(filepath.name, full_content)
        }
    finally:
        pdf.close()


def content_hash(content: str) -> str:
    """Generate short hash for content deduplication."""
    return hashlib.md5(content.encode()).hexdigest()[:8]


def build_catalog(source_dir: Path, output_dir: Path) -> dict:
    """Build complete slide catalog from source directory."""
    catalog = {
        "version": "1.0",
        "generated": datetime.now().isoformat(),
        "source_path": str(source_dir),
        "statistics": {
            "total_presentations": 0,
            "total_slides": 0,
            "pptx_count": 0,
            "pdf_count": 0,
            "skipped_count": 0,
            "error_count": 0,
            "date_range": {"earliest": "", "latest": ""}
        },
        "presentations": []
    }

    # Find all files
    files = list(source_dir.glob("**/*"))
    pptx_files = [f for f in files if f.suffix.lower() == '.pptx']
    pdf_files = [f for f in files if f.suffix.lower() == '.pdf']
    key_files = [f for f in files if f.suffix.lower() == '.key']

    print(f"Found: {len(pptx_files)} PPTX, {len(pdf_files)} PDF, {len(key_files)} KEY (skipping)")

    pres_id = 0
    dates_found = []
    content_hashes = {}  # For dedup detection

    # Process PPTX files
    print("\nProcessing PPTX files...")
    for filepath in sorted(pptx_files):
        pres_id += 1
        print(f"  [{pres_id}] {filepath.name}")

        result = extract_pptx(filepath)
        if not result:
            catalog["statistics"]["error_count"] += 1
            continue

        # Get relative path from Data root
        rel_path = filepath.relative_to(source_dir.parent.parent)
        date_str = parse_date_from_filename(filepath.name)
        if date_str:
            dates_found.append(date_str)

        # Check for content duplicates
        full_hash = content_hash(result["full_content"])
        is_duplicate = full_hash in content_hashes
        if is_duplicate:
            duplicate_of = content_hashes[full_hash]
        else:
            content_hashes[full_hash] = f"pres-{pres_id:03d}"
            duplicate_of = None

        entry = {
            "id": f"pres-{pres_id:03d}",
            "source_file": str(rel_path),
            "format": "pptx",
            "title": filepath.stem.replace('_', ' ').replace('-', ' '),
            "created": date_str,
            "audience": result["audience"],
            "slide_count": result["slide_count"],
            "slides": result["slides"],
            "content_hash": full_hash
        }

        if duplicate_of:
            entry["duplicate_of"] = duplicate_of

        catalog["presentations"].append(entry)
        catalog["statistics"]["pptx_count"] += 1
        catalog["statistics"]["total_slides"] += result["slide_count"]

    # Check which PDFs have PPTX equivalents
    pptx_bases = {f.stem.lower() for f in pptx_files}

    # Process PDF files
    print("\nProcessing PDF files...")
    for filepath in sorted(pdf_files):
        pres_id += 1
        print(f"  [{pres_id}] {filepath.name}")

        has_pptx = filepath.stem.lower() in pptx_bases

        result = extract_pdf(filepath)
        if not result:
            catalog["statistics"]["error_count"] += 1
            continue

        rel_path = filepath.relative_to(source_dir.parent.parent)
        date_str = parse_date_from_filename(filepath.name)
        if date_str:
            dates_found.append(date_str)

        # Check for content duplicates
        full_hash = content_hash(result["full_content"])
        is_duplicate = full_hash in content_hashes
        if is_duplicate:
            duplicate_of = content_hashes[full_hash]
        else:
            content_hashes[full_hash] = f"pres-{pres_id:03d}"
            duplicate_of = None

        entry = {
            "id": f"pres-{pres_id:03d}",
            "source_file": str(rel_path),
            "format": "pdf",
            "title": filepath.stem.replace('_', ' ').replace('-', ' '),
            "created": date_str,
            "audience": result["audience"],
            "slide_count": result["slide_count"],
            "slides": result["slides"],
            "has_pptx_original": has_pptx,
            "content_hash": full_hash
        }

        if duplicate_of:
            entry["duplicate_of"] = duplicate_of

        catalog["presentations"].append(entry)
        catalog["statistics"]["pdf_count"] += 1
        catalog["statistics"]["total_slides"] += result["slide_count"]

    # Update statistics
    catalog["statistics"]["total_presentations"] = len(catalog["presentations"])
    catalog["statistics"]["skipped_count"] = len(key_files)

    if dates_found:
        dates_found.sort()
        catalog["statistics"]["date_range"]["earliest"] = dates_found[0]
        catalog["statistics"]["date_range"]["latest"] = dates_found[-1]

    return catalog


def generate_report(catalog: dict, output_dir: Path) -> str:
    """Generate markdown report from catalog."""
    stats = catalog["statistics"]

    # Collect all keywords
    all_keywords = []
    audience_counts = Counter()

    for pres in catalog["presentations"]:
        audience_counts[pres["audience"]] += 1
        for slide in pres["slides"]:
            all_keywords.extend(slide.get("keywords", []))

    keyword_counts = Counter(all_keywords)
    top_keywords = keyword_counts.most_common(30)

    # Find duplicates
    duplicates = [p for p in catalog["presentations"] if "duplicate_of" in p]

    report = f"""# Organization Slide Catalog Report

Generated: {catalog["generated"][:10]}

## Summary

| Metric | Count |
|--------|-------|
| Total Presentations | {stats["total_presentations"]} |
| Total Slides | {stats["total_slides"]} |
| PPTX Files | {stats["pptx_count"]} |
| PDF Files | {stats["pdf_count"]} |
| Skipped (KEY) | {stats["skipped_count"]} |
| Errors | {stats["error_count"]} |

**Date Range:** {stats["date_range"]["earliest"]} to {stats["date_range"]["latest"]}

## Audience Breakdown

| Audience | Count |
|----------|-------|
"""

    for audience, count in audience_counts.most_common():
        report += f"| {audience} | {count} |\n"

    report += f"""
## Top Keywords

These terms appear most frequently across all slides:

"""

    for keyword, count in top_keywords:
        report += f"- **{keyword}** ({count})\n"

    if duplicates:
        report += f"""
## Duplicate Content Detected

Found {len(duplicates)} presentations with duplicate content:

| Presentation | Duplicate Of |
|--------------|--------------|
"""
        for dup in duplicates:
            report += f"| {dup['title'][:40]} | {dup['duplicate_of']} |\n"

    report += f"""
## Usage

Search the catalog using jq:

```bash
# Find slides mentioning a topic
jq '.presentations[].slides[] | select(.content | test("data sovereignty"; "i"))' slide-catalog.json

# List all investor presentations
jq '.presentations[] | select(.audience == "investor") | .title' slide-catalog.json

# Find presentations from 2018
jq '.presentations[] | select(.created | startswith("2018")) | .title' slide-catalog.json
```

## Source Location

Original files: `{catalog["source_path"]}`

To use a slide with preserved design, open the original file listed in `source_file` field.
"""

    return report


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    source_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    if not source_dir.exists():
        print(f"Error: Source directory not found: {source_dir}")
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Building slide catalog from: {source_dir}")
    print(f"Output directory: {output_dir}")
    print()

    # Build catalog
    catalog = build_catalog(source_dir, output_dir)

    # Write catalog JSON
    catalog_path = output_dir / "slide-catalog.json"
    with open(catalog_path, 'w') as f:
        json.dump(catalog, f, indent=2)
    print(f"\nCatalog written to: {catalog_path}")

    # Generate and write report
    report = generate_report(catalog, output_dir)
    report_path = output_dir / "catalog-report.md"
    with open(report_path, 'w') as f:
        f.write(report)
    print(f"Report written to: {report_path}")

    # Print summary
    stats = catalog["statistics"]
    print(f"\n{'='*50}")
    print(f"COMPLETE: {stats['total_presentations']} presentations, {stats['total_slides']} slides indexed")
    print(f"{'='*50}")


if __name__ == "__main__":
    main()
